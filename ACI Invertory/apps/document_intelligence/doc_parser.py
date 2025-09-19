import hashlib
import magic
import os
import tempfile
import zipfile
import io
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional, Tuple
import logging

import pdfplumber
import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image

from .models import FileMetadata
from .config import settings

logger = logging.getLogger(__name__)

class DocumentParser:
    """Universal document parser supporting multiple file formats with OCR fallback."""
    
    def __init__(self):
        self.supported_types = {
            'application/pdf': self._parse_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._parse_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._parse_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._parse_pptx,
            'text/html': self._parse_html,
            'text/plain': self._parse_text,
            'text/csv': self._parse_text,
            'image/png': self._parse_image,
            'image/jpeg': self._parse_image,
            'image/jpg': self._parse_image,
            'application/zip': self._parse_zip,
        }
    
    def parse_file(self, file_path: str) -> Tuple[str, FileMetadata]:
        """Parse any supported file and return extracted text with metadata."""
        try:
            # Get file metadata
            metadata = self._get_file_metadata(file_path)
            
            # Parse based on MIME type
            parser_func = self.supported_types.get(metadata.mime_type)
            if not parser_func:
                raise ValueError(f"Unsupported file type: {metadata.mime_type}")
            
            text = parser_func(file_path)
            
            if not text.strip():
                logger.warning(f"No text extracted from {file_path}, attempting OCR")
                text = self._ocr_fallback(file_path)
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {str(e)}")
            raise
    
    def _get_file_metadata(self, file_path: str) -> FileMetadata:
        """Extract metadata from file."""
        file_stat = os.stat(file_path)
        
        # Get MIME type
        mime_type = magic.from_file(file_path, mime=True)
        
        # Calculate file hash
        with open(file_path, 'rb') as f:
            file_hash = hashlib.sha256(f.read()).hexdigest()
        
        return FileMetadata(
            file_name=os.path.basename(file_path),
            file_size=file_stat.st_size,
            mime_type=mime_type,
            file_hash=file_hash,
            upload_timestamp=datetime.fromtimestamp(file_stat.st_mtime)
        )
    
    def _parse_pdf(self, file_path: str) -> str:
        """Parse PDF using pdfplumber with PyMuPDF fallback."""
        text = ""
        
        try:
            # Try pdfplumber first
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path}: {e}")
            
        # Fallback to PyMuPDF if no text extracted
        if not text.strip():
            try:
                doc = fitz.open(file_path)
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    text += page.get_text() + "\n"
                doc.close()
            except Exception as e:
                logger.warning(f"PyMuPDF failed for {file_path}: {e}")
        
        return text
    
    def _parse_docx(self, file_path: str) -> str:
        """Parse Word document."""
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    
    def _parse_xlsx(self, file_path: str) -> str:
        """Parse Excel file."""
        workbook = load_workbook(file_path, data_only=True)
        text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
            text += "\n"
        
        return text
    
    def _parse_pptx(self, file_path: str) -> str:
        """Parse PowerPoint presentation."""
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        return text
    
    def _parse_html(self, file_path: str) -> str:
        """Parse HTML file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            return soup.get_text()
    
    def _parse_text(self, file_path: str) -> str:
        """Parse plain text file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _parse_image(self, file_path: str) -> str:
        """Parse image using OCR with multiple preprocessing approaches."""
        try:
            from PIL import ImageEnhance, ImageFilter
            
            image = Image.open(file_path)
            
            # Convert to RGB if needed
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Convert to grayscale for better OCR
            image = image.convert('L')
            
            # Try multiple OCR configurations
            configs = [
                r'--oem 3 --psm 6',  # Default
                r'--oem 3 --psm 7',  # Single text line
                r'--oem 3 --psm 8',  # Single word
                r'--oem 3 --psm 4',  # Single column
                r'--oem 3 --psm 3',  # Fully automatic
            ]
            
            best_text = ""
            best_score = 0
            
            for config in configs:
                try:
                    # Try with current image
                    text = pytesseract.image_to_string(image, config=config)
                    
                    # Score based on recognizable patterns (letters, numbers, spaces)
                    score = sum(1 for c in text if c.isalnum() or c.isspace())
                    
                    if score > best_score:
                        best_score = score
                        best_text = text
                        
                except Exception:
                    continue
            
            # If still poor quality, try with image enhancement
            if best_score < 50:  # Threshold for poor quality
                try:
                    # Enhance contrast
                    enhancer = ImageEnhance.Contrast(image)
                    enhanced = enhancer.enhance(2.0)
                    
                    # Try again with enhanced image
                    enhanced_text = pytesseract.image_to_string(enhanced, config=r'--oem 3 --psm 6')
                    enhanced_score = sum(1 for c in enhanced_text if c.isalnum() or c.isspace())
                    
                    if enhanced_score > best_score:
                        best_text = enhanced_text
                        
                except Exception:
                    pass
            
            return best_text
            
        except Exception as e:
            logger.error(f"OCR failed for image {file_path}: {e}")
            return ""
    
    def _parse_zip(self, file_path: str) -> str:
        """Extract and parse files from ZIP archive."""
        text = ""
        
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            for file_info in zip_file.filelist:
                if file_info.is_dir():
                    continue
                
                try:
                    with zip_file.open(file_info) as extracted_file:
                        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                            temp_file.write(extracted_file.read())
                            temp_path = temp_file.name
                        
                        # Recursively parse extracted file
                        try:
                            extracted_text, _ = self.parse_file(temp_path)
                            text += f"File: {file_info.filename}\n{extracted_text}\n\n"
                        except Exception as e:
                            logger.warning(f"Failed to parse {file_info.filename} from ZIP: {e}")
                        finally:
                            os.unlink(temp_path)
                
                except Exception as e:
                    logger.warning(f"Failed to extract {file_info.filename}: {e}")
        
        return text
    
    def _ocr_fallback(self, file_path: str) -> str:
        """OCR fallback for files with no extractable text."""
        try:
            # For PDFs, convert to images first
            if file_path.lower().endswith('.pdf'):
                return self._ocr_pdf(file_path)
            else:
                # Direct OCR for images
                return self._parse_image(file_path)
        except Exception as e:
            logger.error(f"OCR fallback failed for {file_path}: {e}")
            return ""
    
    def _ocr_pdf(self, file_path: str) -> str:
        """OCR PDF by converting pages to images."""
        text = ""
        try:
            doc = fitz.open(file_path)
            for page_num in range(doc.page_count):
                page = doc[page_num]
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                
                # Create PIL Image from bytes
                image = Image.open(io.BytesIO(img_data))
                page_text = pytesseract.image_to_string(image)
                text += page_text + "\n"
            
            doc.close()
        except Exception as e:
            logger.error(f"PDF OCR failed for {file_path}: {e}")
        
        return text